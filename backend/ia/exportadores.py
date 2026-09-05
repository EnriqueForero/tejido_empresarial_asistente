"""
Descargas del asistente: libro de Excel y presentación de PowerPoint.

Ambos archivos llevan, en un lugar imposible de pasar por alto, la advertencia de
que el contenido lo generó una inteligencia artificial y debe verificarse. Los
dos incluyen además la pregunta, la respuesta y la consulta SQL ejecutada, para
que el resultado sea auditable fuera del aplicativo.

El estilo replica el del resto de descargas del aplicativo: azul noche con acento
ámbar, encabezados congelados, anchos calculados e identificadores como texto.
"""
from __future__ import annotations

import re
import unicodedata
from datetime import datetime
from io import BytesIO
from typing import Any

import xlsxwriter

from backend.config import APP_VERSION, IA_ADVERTENCIA
from backend.ia.forma import clase_de_cifra

NAVY = "#011627"
NAVY_2 = "#062B43"
AMBER = "#FFA400"
AMBER_SOFT = "#FFF4DC"
PAPER = "#F4F7FA"
WHITE = "#FFFFFF"
TEXT = "#0B2233"
MUTED = "#52667A"
BORDER = "#D5DEE5"

FONT_DISPLAY = "Jost"
FONT_BODY = "Maven Pro"

_IDENTIFICADOR = ("NIT", "DIGITO", "DÍGITO", "CODIGO", "CÓDIGO", "COD_", "POSICION", "POSICIÓN", "ID_")
#: Moneda por el nombre de la columna. El asistente inventa sus alias («Total expo
#: 5 anos USD»), así que no se puede depender de la etiqueta exacta del glosario.
#: Sin esto, el libro escribía dólares y pesos con el mismo formato genérico.
_ANCHO_MIN = 11.0
_ANCHO_MAX = 52.0
_MAX_FILAS_EXCEL = 20_000


def _es_identificador(columna: str) -> bool:
    arriba = columna.upper()
    return any(termino in arriba for termino in _IDENTIFICADOR)


def _clase_numerica(columna: str) -> str:
    """«usd», «cop» o «numero»: el formato con el que se escribe una cifra.

    La regla es la de `forma.clase_de_cifra`, para que el archivo, la tabla en
    pantalla y el resumen automático no digan tres cosas distintas del mismo
    número. Aquí un identificador no es un caso aparte: `_es_identificador` ya
    lo desvió a texto antes de llegar a este punto.
    """
    clase = clase_de_cifra(columna)
    return clase if clase in ("usd", "cop") else "numero"


def _ancho(columna: str, valores: list[Any]) -> float:
    largo = len(str(columna))
    for valor in valores[:400]:
        if valor is not None:
            largo = max(largo, len(str(valor)))
    return min(max(largo + 3.0, _ANCHO_MIN), _ANCHO_MAX)


def _nombre_archivo(pregunta: str, extension: str, ahora: datetime | None = None) -> str:
    momento = ahora or datetime.now()
    base = unicodedata.normalize("NFKD", pregunta).encode("ascii", "ignore").decode()
    base = re.sub(r"[^A-Za-z0-9]+", "_", base).strip("_")[:48] or "Consulta"
    return f"ProColombia_Tejido_Empresarial_Asistente_{base}_{momento:%Y%m%d_%H%M}.{extension}"


def nombre_excel(pregunta: str, ahora: datetime | None = None) -> str:
    return _nombre_archivo(pregunta, "xlsx", ahora)


def nombre_pptx(pregunta: str, ahora: datetime | None = None) -> str:
    return _nombre_archivo(pregunta, "pptx", ahora)


# ── Excel ─────────────────────────────────────────────────────────────────


def crear_excel(
    pregunta: str,
    respuesta: str,
    sql: str,
    columnas: list[str],
    filas: list[list[Any]],
    n_filas: int,
    ahora: datetime | None = None,
) -> bytes:
    """Libro con dos hojas: «Respuesta» (contexto y advertencia) y «Datos»."""
    momento = ahora or datetime.now()
    buffer = BytesIO()
    libro = xlsxwriter.Workbook(buffer, {"in_memory": True, "default_date_format": "yyyy-mm-dd"})
    libro.set_properties(
        {
            "title": "Asistente · Tejido Empresarial",
            "author": "ProColombia · Gerencia de Inteligencia Comercial",
            "comments": IA_ADVERTENCIA,
        }
    )

    formatos = {
        "titulo": libro.add_format(
            {"font_name": FONT_DISPLAY, "font_size": 18, "bold": True, "font_color": WHITE, "bg_color": NAVY, "valign": "vcenter", "indent": 1}
        ),
        "subtitulo": libro.add_format(
            {"font_name": FONT_BODY, "font_size": 10, "font_color": "#C9D8E4", "bg_color": NAVY, "valign": "vcenter", "indent": 1}
        ),
        "cinta": libro.add_format({"bg_color": AMBER}),
        "etiqueta": libro.add_format(
            {"font_name": FONT_DISPLAY, "font_size": 10, "bold": True, "font_color": NAVY_2, "valign": "top"}
        ),
        "texto": libro.add_format({"font_name": FONT_BODY, "font_size": 11, "text_wrap": True, "valign": "top", "font_color": TEXT}),
        "aviso": libro.add_format(
            {
                "font_name": FONT_BODY,
                "font_size": 10,
                "text_wrap": True,
                "valign": "top",
                "font_color": "#7A4B00",
                "bg_color": AMBER_SOFT,
                "border": 1,
                "border_color": AMBER,
            }
        ),
        "codigo": libro.add_format(
            {"font_name": "Consolas", "font_size": 9, "text_wrap": True, "valign": "top", "font_color": MUTED, "bg_color": PAPER, "border": 1, "border_color": BORDER}
        ),
        "encabezado": libro.add_format(
            {"font_name": FONT_DISPLAY, "font_size": 10, "bold": True, "font_color": WHITE, "bg_color": NAVY_2, "border": 1, "border_color": NAVY, "text_wrap": True, "valign": "vcenter"}
        ),
        "celda": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "valign": "top"}),
        "celda_par": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "bg_color": PAPER, "valign": "top"}),
        "numero": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "num_format": "#,##0.##"}),
        "numero_par": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "bg_color": PAPER, "num_format": "#,##0.##"}),
        # La misma convención que el Excel estándar: FOB USD con dos decimales,
        # COP sin decimales y en rojo si es negativo.
        "usd": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "num_format": "#,##0.00;[Red]-#,##0.00"}),
        "usd_par": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "bg_color": PAPER, "num_format": "#,##0.00;[Red]-#,##0.00"}),
        "cop": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "num_format": "#,##0;[Red]-#,##0"}),
        "cop_par": libro.add_format({"font_name": FONT_BODY, "font_size": 10, "border": 1, "border_color": BORDER, "bg_color": PAPER, "num_format": "#,##0;[Red]-#,##0"}),
        "id": libro.add_format({"font_name": "Consolas", "font_size": 10, "border": 1, "border_color": BORDER, "align": "left"}),
        "id_par": libro.add_format({"font_name": "Consolas", "font_size": 10, "border": 1, "border_color": BORDER, "bg_color": PAPER, "align": "left"}),
    }

    # ── Hoja 1 · Respuesta ──
    hoja = libro.add_worksheet("Respuesta")
    hoja.hide_gridlines(2)
    hoja.set_column("A:A", 22)
    hoja.set_column("B:B", 96)
    hoja.set_row(0, 34)
    hoja.merge_range(0, 0, 0, 1, "Asistente · Tejido Empresarial", formatos["titulo"])
    hoja.set_row(1, 18)
    hoja.merge_range(1, 0, 1, 1, "ProColombia · Gerencia de Inteligencia Comercial", formatos["subtitulo"])
    hoja.set_row(2, 4)
    hoja.merge_range(2, 0, 2, 1, "", formatos["cinta"])

    fila = 4
    for etiqueta, valor in (
        ("Pregunta", pregunta),
        ("Respuesta", respuesta),
        ("Filas obtenidas", f"{n_filas:,}".replace(",", ".")),
        ("Generado", f"{momento:%Y-%m-%d %H:%M}"),
        ("Versión del aplicativo", APP_VERSION),
    ):
        hoja.write(fila, 0, etiqueta, formatos["etiqueta"])
        hoja.write(fila, 1, valor, formatos["texto"])
        hoja.set_row(fila, 15 + 12 * min(6, len(str(valor)) // 90))
        fila += 1

    fila += 1
    hoja.write(fila, 0, "Advertencia", formatos["etiqueta"])
    hoja.merge_range(fila, 1, fila + 2, 1, IA_ADVERTENCIA, formatos["aviso"])
    fila += 4
    hoja.write(fila, 0, "Consulta ejecutada", formatos["etiqueta"])
    hoja.merge_range(fila, 1, fila + 5, 1, sql or "(sin consulta)", formatos["codigo"])

    # ── Hoja 2 · Datos ──
    if columnas:
        datos = libro.add_worksheet("Datos")
        datos.hide_gridlines(2)
        datos.freeze_panes(1, 0)
        for indice, columna in enumerate(columnas):
            valores = [fila_datos[indice] if indice < len(fila_datos) else None for fila_datos in filas[:400]]
            datos.set_column(indice, indice, _ancho(columna, valores))
            datos.write(0, indice, columna, formatos["encabezado"])
        for numero, fila_datos in enumerate(filas[:_MAX_FILAS_EXCEL], start=1):
            par = numero % 2 == 0
            for indice, columna in enumerate(columnas):
                valor = fila_datos[indice] if indice < len(fila_datos) else None
                if _es_identificador(columna):
                    datos.write_string(numero, indice, "" if valor is None else str(valor), formatos["id_par" if par else "id"])
                elif isinstance(valor, (int, float)) and not isinstance(valor, bool):
                    clase = _clase_numerica(columna)
                    datos.write_number(numero, indice, float(valor), formatos[f"{clase}_par" if par else clase])
                else:
                    datos.write(numero, indice, "" if valor is None else str(valor), formatos["celda_par" if par else "celda"])
        if filas:
            datos.autofilter(0, 0, min(len(filas), _MAX_FILAS_EXCEL), len(columnas) - 1)
        datos.set_landscape()
        datos.repeat_rows(0)

    libro.close()
    return buffer.getvalue()


# ── PowerPoint ────────────────────────────────────────────────────────────

_FILAS_POR_LAMINA = 12
_MAX_LAMINAS_TABLA = 4
_MAX_COLUMNAS_PPTX = 8


def crear_pptx(
    pregunta: str,
    respuesta: str,
    sql: str,
    columnas: list[str],
    filas: list[list[Any]],
    n_filas: int,
    ahora: datetime | None = None,
) -> bytes:
    """Presentación con portada, respuesta, tabla paginada y cierre con la SQL."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    momento = ahora or datetime.now()
    azul = RGBColor.from_string(NAVY.lstrip("#"))
    azul_2 = RGBColor.from_string(NAVY_2.lstrip("#"))
    ambar = RGBColor.from_string(AMBER.lstrip("#"))
    blanco = RGBColor(255, 255, 255)
    gris = RGBColor.from_string(MUTED.lstrip("#"))

    presentacion = Presentation()
    presentacion.slide_width = Inches(13.333)
    presentacion.slide_height = Inches(7.5)
    vacia = presentacion.slide_layouts[6]

    def banda(lamina, texto: str) -> None:
        forma = lamina.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentacion.slide_width, Inches(0.95))
        forma.fill.solid()
        forma.fill.fore_color.rgb = azul
        forma.line.fill.background()
        marco = forma.text_frame
        marco.margin_left = Inches(0.45)
        marco.text = texto
        parrafo = marco.paragraphs[0]
        parrafo.font.size = Pt(20)
        parrafo.font.bold = True
        parrafo.font.color.rgb = blanco
        acento = lamina.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), presentacion.slide_width, Inches(0.06))
        acento.fill.solid()
        acento.fill.fore_color.rgb = ambar
        acento.line.fill.background()

    # Lámina 1 · portada y respuesta
    lamina = presentacion.slides.add_slide(vacia)
    banda(lamina, "Asistente · Tejido Empresarial · ProColombia")
    cuadro = lamina.shapes.add_textbox(Inches(0.6), Inches(1.4), presentacion.slide_width - Inches(1.2), Inches(4.4))
    marco = cuadro.text_frame
    marco.word_wrap = True
    titulo = marco.paragraphs[0]
    titulo.text = pregunta
    titulo.font.size = Pt(22)
    titulo.font.bold = True
    titulo.font.color.rgb = azul_2
    cuerpo = marco.add_paragraph()
    cuerpo.text = respuesta
    cuerpo.font.size = Pt(15)
    cuerpo.space_before = Pt(14)
    pie = marco.add_paragraph()
    pie.text = f"{n_filas:,} fila(s) · Generado {momento:%Y-%m-%d %H:%M} · Aplicativo {APP_VERSION}".replace(",", ".")
    pie.font.size = Pt(11)
    pie.font.color.rgb = gris
    pie.space_before = Pt(18)

    aviso = lamina.shapes.add_textbox(Inches(0.6), Inches(6.0), presentacion.slide_width - Inches(1.2), Inches(1.1))
    marco_aviso = aviso.text_frame
    marco_aviso.word_wrap = True
    parrafo_aviso = marco_aviso.paragraphs[0]
    parrafo_aviso.text = IA_ADVERTENCIA
    parrafo_aviso.font.size = Pt(10)
    parrafo_aviso.font.italic = True
    parrafo_aviso.font.color.rgb = gris

    # Láminas de tabla
    if columnas and filas:
        visibles = columnas[:_MAX_COLUMNAS_PPTX]
        tope = _FILAS_POR_LAMINA * _MAX_LAMINAS_TABLA
        for comienzo in range(0, min(len(filas), tope), _FILAS_POR_LAMINA):
            trozo = filas[comienzo : comienzo + _FILAS_POR_LAMINA]
            lamina = presentacion.slides.add_slide(vacia)
            banda(lamina, f"Resultados ({comienzo + 1}–{comienzo + len(trozo)} de {n_filas:,})".replace(",", "."))
            forma = lamina.shapes.add_table(
                rows=len(trozo) + 1,
                cols=len(visibles),
                left=Inches(0.45),
                top=Inches(1.35),
                width=presentacion.slide_width - Inches(0.9),
                height=Inches(0.32) * (len(trozo) + 1),
            )
            tabla = forma.table
            for columna_indice, columna in enumerate(visibles):
                celda = tabla.cell(0, columna_indice)
                celda.text = str(columna)
                celda.fill.solid()
                celda.fill.fore_color.rgb = azul_2
                parrafo = celda.text_frame.paragraphs[0]
                parrafo.font.color.rgb = blanco
                parrafo.font.size = Pt(11)
                parrafo.font.bold = True
            for fila_indice, fila_datos in enumerate(trozo, start=1):
                for columna_indice in range(len(visibles)):
                    valor = fila_datos[columna_indice] if columna_indice < len(fila_datos) else ""
                    if isinstance(valor, float):
                        valor = f"{valor:,.2f}".replace(",", "·").replace(".", ",").replace("·", ".")
                    elif isinstance(valor, int) and not isinstance(valor, bool):
                        valor = f"{valor:,}".replace(",", ".")
                    celda = tabla.cell(fila_indice, columna_indice)
                    celda.text = "" if valor is None else str(valor)
                    celda.text_frame.paragraphs[0].font.size = Pt(10)
        if len(columnas) > _MAX_COLUMNAS_PPTX or len(filas) > tope:
            nota = presentacion.slides[-1].shapes.add_textbox(
                Inches(0.45), Inches(6.6), presentacion.slide_width - Inches(0.9), Inches(0.5)
            )
            parrafo = nota.text_frame.paragraphs[0]
            parrafo.text = "La presentación resume el resultado; el archivo de Excel trae todas las filas y columnas."
            parrafo.font.size = Pt(10)
            parrafo.font.color.rgb = gris

    # Lámina de cierre · trazabilidad
    lamina = presentacion.slides.add_slide(vacia)
    banda(lamina, "Trazabilidad")
    cuadro = lamina.shapes.add_textbox(Inches(0.6), Inches(1.4), presentacion.slide_width - Inches(1.2), Inches(5.4))
    marco = cuadro.text_frame
    marco.word_wrap = True
    encabezado = marco.paragraphs[0]
    encabezado.text = "Consulta ejecutada en Snowflake"
    encabezado.font.size = Pt(15)
    encabezado.font.bold = True
    encabezado.font.color.rgb = azul_2
    codigo = marco.add_paragraph()
    codigo.text = (sql or "(sin consulta)")[:2200]
    codigo.font.size = Pt(9)
    codigo.font.name = "Consolas"
    final = marco.add_paragraph()
    final.text = IA_ADVERTENCIA
    final.font.size = Pt(10)
    final.font.italic = True
    final.font.color.rgb = gris
    final.space_before = Pt(16)

    buffer = BytesIO()
    presentacion.save(buffer)
    return buffer.getvalue()
